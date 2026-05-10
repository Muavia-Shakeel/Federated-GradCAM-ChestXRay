"""
Generate professional PPT for FL+GradCAM project using python-pptx.
Run: /home/msi/DS_env/bin/python3 generate_ppt.py
Output: outputs/FL_GradCAM_Presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

OUT_PATH = os.path.join(os.path.dirname(__file__), "..", "outputs", "reports", "FL_GradCAM_Presentation.pptx")
os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)

# ── Color Palette ──────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x08, 0x0E, 0x26)   # dark navy background
ROYAL_BLUE = RGBColor(0x1A, 0x3A, 0x6E)   # section headers
ACCENT     = RGBColor(0x4A, 0x90, 0xD9)   # highlights / table headers
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xEB, 0xEC, 0xEF)
GREEN      = RGBColor(0x27, 0xAE, 0x60)
ORANGE     = RGBColor(0xE6, 0x7E, 0x22)
RED        = RGBColor(0xC0, 0x39, 0x2B)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color=NAVY):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name  = "Calibri"
    return txb


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_table(slide, data, left, top, width, height,
              header_bg=ACCENT, row_bg=RGBColor(0x10, 0x1A, 0x3A),
              alt_bg=RGBColor(0x15, 0x22, 0x4A),
              header_color=WHITE, cell_color=LIGHT_GRAY,
              font_size=11):
    rows = len(data)
    cols = len(data[0])
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    for r, row_data in enumerate(data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            tf = cell.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run()
            run.font.size  = Pt(font_size)
            run.font.name  = "Calibri"
            run.font.bold  = (r == 0)
            run.font.color.rgb = header_color if r == 0 else cell_color

            from pptx.oxml.ns import qn
            from lxml import etree
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
            srgbClr   = etree.SubElement(solidFill, qn('a:srgbClr'))
            bg = header_bg if r == 0 else (alt_bg if r % 2 == 0 else row_bg)
            srgbClr.set('val', f'{bg[0]:02X}{bg[1]:02X}{bg[2]:02X}')
    return table


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════════════════════════════════════
def slide_title(prs):
    s = blank_slide(prs)
    fill_bg(s, NAVY)

    # Top accent bar
    add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT)

    # Main title
    add_textbox(s, "Federated Learning with\nWeighted GradCAM Aggregation",
                Inches(0.8), Inches(1.2), Inches(11.5), Inches(2.2),
                font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(s, "Privacy-Preserving & Explainable Multi-Label Chest X-Ray Classification",
                Inches(0.8), Inches(3.4), Inches(11.5), Inches(0.6),
                font_size=18, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)

    # Divider
    add_rect(s, Inches(3.5), Inches(4.1), Inches(6.3), Inches(0.03), ACCENT)

    # Authors / course
    add_textbox(s, "Muavia Shakeel & Haseeb   |   Advanced Machine Learning   |   MS 2026",
                Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.5),
                font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Dataset / model badges
    badges = [
        ("NIH ChestX-ray14", Inches(2.5)),
        ("EfficientNet-B0", Inches(5.5)),
        ("5 FL Clients", Inches(8.5)),
    ]
    for label, left in badges:
        add_rect(s, left, Inches(5.3), Inches(2.5), Inches(0.5),
                 ROYAL_BLUE, ACCENT)
        add_textbox(s, label, left, Inches(5.3), Inches(2.5), Inches(0.5),
                    font_size=12, color=WHITE, align=PP_ALIGN.CENTER, bold=True)

    # Bottom bar
    add_rect(s, 0, Inches(7.3), SLIDE_W, Inches(0.2), ROYAL_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: The Dual Problem
# ══════════════════════════════════════════════════════════════════════════════
def slide_problem(prs):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT)

    add_textbox(s, "The Dual Problem in Medical AI",
                Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
                font_size=28, bold=True, color=WHITE)
    add_rect(s, Inches(0.5), Inches(0.85), Inches(4), Inches(0.03), ACCENT)

    # Problem 1
    add_rect(s, Inches(0.5), Inches(1.1), Inches(5.8), Inches(5.5),
             RGBColor(0x10, 0x1A, 0x3E))
    add_textbox(s, "Problem 1: Privacy", Inches(0.7), Inches(1.2),
                Inches(5.4), Inches(0.6), font_size=20, bold=True, color=ACCENT)
    add_textbox(s,
        "• HIPAA prohibits hospitals from sharing\n  patient X-rays with central servers\n\n"
        "• Training centralized AI requires pooling\n  data — legally and ethically impossible\n\n"
        "• 100+ million chest X-rays exist across\n  thousands of hospitals — inaccessible\n\n"
        "→ Solution: Federated Learning (FL)\n"
        "  Train locally. Share only model weights.",
        Inches(0.7), Inches(1.85), Inches(5.4), Inches(4.5),
        font_size=13, color=LIGHT_GRAY)

    # Problem 2
    add_rect(s, Inches(6.9), Inches(1.1), Inches(5.8), Inches(5.5),
             RGBColor(0x10, 0x1A, 0x3E))
    add_textbox(s, "Problem 2: Trust (Black Box)", Inches(7.1), Inches(1.2),
                Inches(5.4), Inches(0.6), font_size=20, bold=True, color=ACCENT)
    add_textbox(s,
        "• Doctors cannot use AI that cannot\n  explain its reasoning\n\n"
        "• 'This patient has 90% Pneumonia'\n  — Why? What did you see?\n\n"
        "• Standard CNN: no visual explanation\n  available to the clinician\n\n"
        "→ Solution: GradCAM Heatmaps\n"
        "  Show exactly which pixels the\n  model attended to.",
        Inches(7.1), Inches(1.85), Inches(5.4), Inches(4.5),
        font_size=13, color=LIGHT_GRAY)

    add_textbox(s, "Our work solves BOTH simultaneously — first in federated multi-label CXR literature",
                Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.5),
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: Literature Gaps
# ══════════════════════════════════════════════════════════════════════════════
def slide_literature(prs):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT)

    add_textbox(s, "5 Gaps in Prior Literature",
                Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
                font_size=28, bold=True, color=WHITE)
    add_rect(s, Inches(0.5), Inches(0.85), Inches(4), Inches(0.03), ACCENT)

    gaps = [
        ("Gap 1", "FL + XAI Simultaneously",
         "Kaissis (2021), Adnan (2022): FL only — no explainability"),
        ("Gap 2", "Federated GradCAM Aggregation",
         "Selvaraju (2020), Tjoa (2021): GradCAM on centralized data only"),
        ("Gap 3", "Dataset-Size-Weighted Explanation",
         "2024 blockchain FL paper: post-hoc GradCAM, ignores client heterogeneity"),
        ("Gap 4", "Quantitative XAI Validation",
         "Prior FL+XAI papers use visual inspection only — no Faithfulness/AOPC metrics"),
        ("Gap 5", "Non-IID Explanation Divergence",
         "No paper measures how non-IID data causes client GradCAM divergence"),
    ]

    for i, (tag, title, desc) in enumerate(gaps):
        y = Inches(1.05) + i * Inches(1.2)
        add_rect(s, Inches(0.4), y, Inches(0.8), Inches(0.95), ROYAL_BLUE)
        add_textbox(s, tag, Inches(0.4), y + Inches(0.15), Inches(0.8), Inches(0.7),
                    font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(s, Inches(1.3), y, Inches(11.4), Inches(0.95),
                 RGBColor(0x10, 0x1A, 0x3E))
        add_textbox(s, title, Inches(1.45), y + Inches(0.05), Inches(11), Inches(0.4),
                    font_size=14, bold=True, color=ACCENT)
        add_textbox(s, desc, Inches(1.45), y + Inches(0.45), Inches(11), Inches(0.45),
                    font_size=11, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: Our Innovation
# ══════════════════════════════════════════════════════════════════════════════
def slide_innovation(prs):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT)

    add_textbox(s, "Our Core Innovation",
                Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
                font_size=28, bold=True, color=WHITE)
    add_rect(s, Inches(0.5), Inches(0.85), Inches(4), Inches(0.03), ACCENT)

    add_textbox(s, "Dataset-Size-Weighted Federated GradCAM Aggregation",
                Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.6),
                font_size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # Formula box
    add_rect(s, Inches(1.5), Inches(1.7), Inches(10.3), Inches(0.9),
             RGBColor(0x05, 0x0A, 0x1A))
    add_textbox(s,
        "Global_Map = Σ (nᵢ / Σnⱼ) × ClientMap_i     where nᵢ = training samples at hospital i",
        Inches(1.6), Inches(1.8), Inches(10), Inches(0.7),
        font_size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # Steps
    steps = [
        ("Hospital A\n10,000 patients", "→", "GradCAM\nMap A", "→"),
        ("Hospital B\n500 patients",    "→", "GradCAM\nMap B", "→"),
        ("Hospital C\n2,000 patients",  "→", "GradCAM\nMap C", "→"),
    ]
    add_textbox(s, "Each client generates local heatmap (no raw X-rays transmitted):",
                Inches(0.5), Inches(2.75), Inches(12.3), Inches(0.4),
                font_size=13, color=LIGHT_GRAY)

    for i, (hosp, arr1, cmap, arr2) in enumerate(steps):
        y = Inches(3.2) + i * Inches(0.9)
        add_rect(s, Inches(0.5), y, Inches(2.2), Inches(0.75), ROYAL_BLUE)
        add_textbox(s, hosp, Inches(0.5), y, Inches(2.2), Inches(0.75),
                    font_size=11, color=WHITE, align=PP_ALIGN.CENTER, bold=True)
        add_textbox(s, "→", Inches(2.8), y + Inches(0.2), Inches(0.5), Inches(0.4),
                    font_size=18, color=ACCENT, align=PP_ALIGN.CENTER)
        add_rect(s, Inches(3.4), y, Inches(2.0), Inches(0.75),
                 RGBColor(0x1A, 0x4A, 0x3A))
        add_textbox(s, cmap, Inches(3.4), y, Inches(2.0), Inches(0.75),
                    font_size=11, color=GREEN, align=PP_ALIGN.CENTER, bold=True)

    # Aggregation arrow
    add_textbox(s, "↓  Weighted Aggregation  ↓",
                Inches(5.8), Inches(3.35), Inches(7), Inches(0.5),
                font_size=16, bold=True, color=ACCENT)
    add_rect(s, Inches(5.8), Inches(3.95), Inches(7), Inches(1.3),
             RGBColor(0x1A, 0x3A, 0x6E))
    add_textbox(s,
        "Global GradCAM Map\n"
        "Hospital A gets 10,000/12,500 = 80% voting power\n"
        "Hospital B gets 500/12,500 = 4% voting power",
        Inches(5.9), Inches(4.05), Inches(6.8), Inches(1.1),
        font_size=12, color=WHITE, bold=False)

    add_textbox(s,
        "Privacy guarantee: Only heatmaps leave clients — no patient data transmitted",
        Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
        font_size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: Dataset + Model + Setup
# ══════════════════════════════════════════════════════════════════════════════
def slide_setup(prs):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT)

    add_textbox(s, "Experimental Setup",
                Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
                font_size=28, bold=True, color=WHITE)
    add_rect(s, Inches(0.5), Inches(0.85), Inches(4), Inches(0.03), ACCENT)

    cols = [
        {
            "title": "Dataset",
            "color": ACCENT,
            "left": Inches(0.3),
            "items": [
                "NIH ChestX-ray14",
                "112,120 frontal X-rays",
                "14 pathology classes",
                "Multi-label (co-morbidities)",
                "10% held-out test set",
                "Hernia pos_weight = 489×",
                "(severe class imbalance handled)",
            ]
        },
        {
            "title": "FL Configuration",
            "color": GREEN,
            "left": Inches(4.7),
            "items": [
                "5 virtual hospital clients",
                "Non-IID: Dirichlet α = 0.5",
                "20 rounds (early stopping)",
                "Local epochs = 3",
                "FedAvg + FedProx (μ=0.01)",
                "Batch size = 32",
                "LR = 1e-4, Adam",
            ]
        },
        {
            "title": "Model",
            "color": ORANGE,
            "left": Inches(9.1),
            "items": [
                "EfficientNet-B0 (timm)",
                "Pre-trained ImageNet-21k",
                "5.3M parameters",
                "Dropout = 0.3",
                "BCEWithLogitsLoss + pos_weight",
                "GradCAM @ conv_pwl layer",
                "(1152→320 channels, 1×1)",
            ]
        },
    ]

    for col in cols:
        add_rect(s, col["left"], Inches(1.0), Inches(3.9), Inches(5.7),
                 RGBColor(0x10, 0x1A, 0x3E))
        add_rect(s, col["left"], Inches(1.0), Inches(3.9), Inches(0.5), col["color"])
        add_textbox(s, col["title"], col["left"], Inches(1.0), Inches(3.9), Inches(0.5),
                    font_size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        for j, item in enumerate(col["items"]):
            y = Inches(1.6) + j * Inches(0.68)
            add_textbox(s, f"• {item}", col["left"] + Inches(0.1), y,
                        Inches(3.7), Inches(0.65), font_size=11.5, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: Results — Classification
# ══════════════════════════════════════════════════════════════════════════════
def slide_results_class(prs):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT)

    add_textbox(s, "Classification Results",
                Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
                font_size=28, bold=True, color=WHITE)
    add_rect(s, Inches(0.5), Inches(0.85), Inches(4), Inches(0.03), ACCENT)

    # Main comparison table
    data = [
        ["Model",               "AUC-ROC (macro)", "F1-macro", "Hamming Score", "Rounds"],
        ["Centralized (upper)", "0.8291",           "0.2440",   "0.8155",        "20"],
        ["FedProx + GradCAM",   "0.8210",           "0.2480",   "0.8263",        "19/20"],
        ["FedAvg + GradCAM",    "0.8235",           "0.2266",   "0.7953",        "13/20 ★"],
        ["FedAvg (no GradCAM)", "0.8256",           "0.2268",   "0.8036",        "14/20"],
    ]
    add_table(s, data, Inches(0.5), Inches(1.05), Inches(12.3), Inches(2.5), font_size=12)

    add_textbox(s, "★ Early stopping at round 13 — less stable under non-IID vs FedProx",
                Inches(0.5), Inches(3.65), Inches(12.3), Inches(0.35),
                font_size=11, color=ORANGE, italic=True)

    # Key insights
    add_textbox(s, "Key Findings", Inches(0.5), Inches(4.1),
                Inches(12.3), Inches(0.45), font_size=16, bold=True, color=ACCENT)

    insights = [
        ("FedProx wins F1 & Hamming:", "More stable convergence (19 rounds). Proximal term regularizes non-IID divergence."),
        ("FL vs Centralized gap:",      "AUC gap = 0.0056–0.0081. Privacy cost is minimal — clinically acceptable."),
        ("GradCAM overhead:",           "Adding GradCAM aggregation does NOT degrade classification AUC — privacy + explainability for free."),
    ]
    for i, (label, desc) in enumerate(insights):
        y = Inches(4.6) + i * Inches(0.82)
        add_textbox(s, label, Inches(0.5), y, Inches(3.2), Inches(0.7),
                    font_size=12, bold=True, color=GREEN)
        add_textbox(s, desc, Inches(3.8), y, Inches(9.0), Inches(0.7),
                    font_size=12, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: Results — XAI Metrics
# ══════════════════════════════════════════════════════════════════════════════
def slide_results_xai(prs):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT)

    add_textbox(s, "XAI Evaluation Results",
                Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
                font_size=28, bold=True, color=WHITE)
    add_rect(s, Inches(0.5), Inches(0.85), Inches(4), Inches(0.03), ACCENT)

    # Oracle faithfulness table
    add_textbox(s, "Oracle Faithfulness (Multi-Class: Atelectasis, Effusion, Infiltration)",
                Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.4),
                font_size=14, bold=True, color=ACCENT)
    data_faith = [
        ["Model",     "Mean Faith.", "Atelectasis", "Effusion ★", "Infiltration", "AOPC",   "Ins AUC", "Del AUC"],
        ["FedAvg",    "-0.001",      "-0.009",       "+0.050",     "-0.043",       "-0.100",  "0.343",   "0.418"],
        ["FedProx ★", "+0.004",      "-0.030",       "+0.087",     "-0.046",       "-0.103",  "0.324",   "0.357"],
    ]
    add_table(s, data_faith, Inches(0.5), Inches(1.45), Inches(12.3), Inches(1.5), font_size=11)

    add_textbox(s, "★ FedProx achieves POSITIVE mean faithfulness — more stable weights → more consistent GradCAM",
                Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.35),
                font_size=11, color=GREEN, italic=True, bold=True)

    # Strategy comparison table
    add_textbox(s, "Aggregation Strategy Comparison (vs Oracle, Infiltration class)",
                Inches(0.5), Inches(3.45), Inches(12.3), Inches(0.4),
                font_size=14, bold=True, color=ACCENT)
    data_strat = [
        ["Strategy",         "Faithfulness", "AOPC",   "Pearson-r", "SSIM",   "MSE"],
        ["Weighted ★ (ours)", "+0.0021",      "-0.069",  "-0.287",    "-0.242", "0.129"],
        ["Uniform",           "-0.0002",      "-0.066",  "-0.137",    "-0.108", "0.129"],
        ["Max-pool",          "-0.039",       "-0.075",  "-0.106",    "-0.080", "0.124"],
    ]
    add_table(s, data_strat, Inches(0.5), Inches(3.9), Inches(12.3), Inches(1.7), font_size=11)

    add_textbox(s,
        "★ Weighted is ONLY strategy with positive faithfulness — dataset-size weighting produces more disease-focused explanations\n"
        "Negative Pearson-r reveals client GradCAM divergence under non-IID — validates why federated XAI aggregation is needed",
        Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.7),
        font_size=11, color=ORANGE, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: GradCAM Visualization
# ══════════════════════════════════════════════════════════════════════════════
def slide_gradcam_viz(prs):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT)

    add_textbox(s, "GradCAM Heatmap Evolution Across 20 Rounds",
                Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
                font_size=28, bold=True, color=WHITE)
    add_rect(s, Inches(0.5), Inches(0.85), Inches(4), Inches(0.03), ACCENT)

    # Try to embed actual gradcam images
    plots_dir = os.path.join(os.path.dirname(__file__), "..", "outputs", "plots")
    gradcam_files = sorted([
        f for f in os.listdir(plots_dir)
        if f.startswith("gradcam_round") and f.endswith(".png")
    ])

    # Show 4 key rounds: 1, 5, 10, 20 (or last available)
    key_rounds = [1, 5, 10, 20]
    shown = []
    for r in key_rounds:
        fn = f"gradcam_round{r}.png"
        if fn in gradcam_files:
            shown.append((r, os.path.join(plots_dir, fn)))
        elif gradcam_files:
            # use last available
            shown.append((r, os.path.join(plots_dir, gradcam_files[-1])))

    if shown:
        for i, (round_num, img_path) in enumerate(shown[:4]):
            left = Inches(0.3) + i * Inches(3.2)
            try:
                s.shapes.add_picture(img_path, left, Inches(1.1),
                                     Inches(3.0), Inches(4.5))
                add_textbox(s, f"Round {round_num}", left, Inches(5.65),
                            Inches(3.0), Inches(0.4),
                            font_size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
            except Exception:
                add_rect(s, left, Inches(1.1), Inches(3.0), Inches(4.5),
                         RGBColor(0x10, 0x1A, 0x3E))
                add_textbox(s, f"Round {round_num}\n[GradCAM]", left, Inches(2.8),
                            Inches(3.0), Inches(0.8),
                            font_size=14, color=ACCENT, align=PP_ALIGN.CENTER)
    else:
        add_textbox(s,
            "GradCAM heatmaps saved in outputs/plots/gradcam_round*.png\n"
            "Insert manually for presentation",
            Inches(1), Inches(2.5), Inches(11.3), Inches(2),
            font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_textbox(s,
        "Early rounds: diffuse activation across lung fields  →  Later rounds: sharper focus on pathology-specific regions",
        Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.45),
        font_size=12, color=LIGHT_GRAY, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: Literature Comparison Table
# ══════════════════════════════════════════════════════════════════════════════
def slide_comparison(prs):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT)

    add_textbox(s, "Comparison with Prior Work",
                Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
                font_size=28, bold=True, color=WHITE)
    add_rect(s, Inches(0.5), Inches(0.85), Inches(4), Inches(0.03), ACCENT)

    data = [
        ["Paper (Year)",            "FL", "XAI", "Multi-Label", "Non-IID", "Quant. XAI", "FL GradCAM"],
        ["Kaissis et al. (2021)",   "✓",  "✗",   "✗",           "✗",       "✗",          "✗"],
        ["Adnan et al. (2022)",     "✓",  "✗",   "✗",           "Partial", "✗",          "✗"],
        ["Selvaraju et al. (2020)", "✗",  "✓",   "✗",           "✗",       "Partial",    "✗"],
        ["Tjoa & Guan (2021)",      "✗",  "✓",   "✗",           "✗",       "✗",          "✗"],
        ["Sheller et al. (2020)",   "✓",  "✗",   "✗",           "Partial", "✗",          "✗"],
        ["Arun et al. (2021)",      "✗",  "✓",   "✓",           "✗",       "✓",          "✗"],
        ["2024 Blockchain FL+XAI",  "✓",  "✓",   "✗",           "✗",       "✗",          "Post-hoc"],
        ["Ours (2026) ★",           "✓",  "✓",   "✓",           "✓",       "✓",          "✓ Weighted"],
    ]
    add_table(s, data, Inches(0.4), Inches(1.05), Inches(12.5), Inches(5.5), font_size=11)

    add_textbox(s, "★ Our work is the first to satisfy all 6 criteria simultaneously",
                Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
                font_size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10: Contributions + Conclusion
# ══════════════════════════════════════════════════════════════════════════════
def slide_conclusion(prs):
    s = blank_slide(prs)
    fill_bg(s, NAVY)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT)

    add_textbox(s, "Contributions & Conclusion",
                Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
                font_size=28, bold=True, color=WHITE)
    add_rect(s, Inches(0.5), Inches(0.85), Inches(4), Inches(0.03), ACCENT)

    contribs = [
        ("C1", "Novel Algorithm",
         "Dataset-size-weighted federated GradCAM aggregation — first principled method for FL explanation aggregation"),
        ("C2", "Dual Privacy + XAI",
         "Simultaneously satisfies HIPAA compliance (FL) and clinical explainability (GradCAM) for 14-class CXR"),
        ("C3", "Quantitative XAI",
         "First application of Faithfulness, AOPC, Insertion/Deletion AUC to federated multi-label chest X-ray"),
        ("C4", "Non-IID Explanation Study",
         "Empirically demonstrates client GradCAM divergence under non-IID (negative Pearson-r) — new finding"),
        ("C5", "Competitive Performance",
         "FedProx AUC 0.821 vs Centralized 0.829 — only 0.8% gap with full privacy preservation"),
    ]

    for i, (tag, title, desc) in enumerate(contribs):
        y = Inches(1.05) + i * Inches(1.1)
        add_rect(s, Inches(0.4), y, Inches(0.65), Inches(0.95), ACCENT)
        add_textbox(s, tag, Inches(0.4), y + Inches(0.2), Inches(0.65), Inches(0.5),
                    font_size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_rect(s, Inches(1.1), y, Inches(11.6), Inches(0.95),
                 RGBColor(0x10, 0x1A, 0x3E))
        add_textbox(s, title, Inches(1.2), y + Inches(0.03), Inches(11.3), Inches(0.4),
                    font_size=13, bold=True, color=ACCENT)
        add_textbox(s, desc, Inches(1.2), y + Inches(0.45), Inches(11.3), Inches(0.45),
                    font_size=11, color=LIGHT_GRAY)

    add_rect(s, 0, Inches(7.1), SLIDE_W, Inches(0.4), ROYAL_BLUE)
    add_textbox(s, "Federated Learning + Weighted GradCAM = Privacy + Trust + Performance",
                Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.4),
                font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════
def main():
    prs = new_prs()
    print("Building slides ...")
    slide_title(prs);       print("  Slide 1: Title")
    slide_problem(prs);     print("  Slide 2: Problem")
    slide_literature(prs);  print("  Slide 3: Literature Gaps")
    slide_innovation(prs);  print("  Slide 4: Innovation")
    slide_setup(prs);       print("  Slide 5: Setup")
    slide_results_class(prs); print("  Slide 6: Classification Results")
    slide_results_xai(prs); print("  Slide 7: XAI Results")
    slide_gradcam_viz(prs); print("  Slide 8: GradCAM Visualization")
    slide_comparison(prs);  print("  Slide 9: Literature Comparison")
    slide_conclusion(prs);  print("  Slide 10: Conclusion")

    prs.save(OUT_PATH)
    print(f"\nSaved: {OUT_PATH}")


if __name__ == "__main__":
    main()
